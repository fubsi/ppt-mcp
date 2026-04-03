from __future__ import annotations

import math
from io import BytesIO
from pathlib import Path
from typing import Optional
from urllib.error import HTTPError, URLError
from urllib.parse import urlparse
from urllib.request import urlopen

from pptx.oxml.ns import qn


EMU_PER_POINT = 12700
DEFAULT_FONT_SIZE_PT = 18.0
AVG_CHAR_WIDTH_FACTOR = 0.52
LINE_HEIGHT_FACTOR = 1.2


def get_slide_by_id(pptx_object, slide_id: int):
	"""Return the slide object matching slide_id or None."""
	for slide in pptx_object.slides:
		if slide.slide_id == slide_id:
			return slide
	return None


def get_slide_with_index_by_id(pptx_object, slide_id: int):
	"""Return (slide, index) for slide_id or (None, None) when missing."""
	for index, slide in enumerate(pptx_object.slides):
		if slide.slide_id == slide_id:
			return slide, index
	return None, None


def get_placeholder_by_shape_id(slide, placeholder_shape_id: int):
	"""Return placeholder object matching placeholder_shape_id or None."""
	for placeholder in slide.placeholders:
		if placeholder.shape_id == placeholder_shape_id:
			return placeholder
	return None


def resolve_picture_source(image_source: str):
	"""Resolve local/URL image source to (source, source_label, error_dict)."""
	parsed = urlparse(image_source)
	if parsed.scheme in ("http", "https"):
		try:
			with urlopen(image_source, timeout=10) as response:
				image_data = response.read()
			if not image_data:
				return None, None, {"error": "Image URL returned no data."}
			return BytesIO(image_data), image_source, None
		except (HTTPError, URLError, TimeoutError, ValueError) as exc:
			return None, None, {
				"error": "Failed to download image from URL.",
				"details": str(exc),
			}

	resolved_path = Path(image_source).expanduser().resolve()
	if not resolved_path.exists() or not resolved_path.is_file():
		return None, None, {"error": "Image file not found."}

	return str(resolved_path), str(resolved_path), None


def remove_shapes_by_ids(
	presentations: dict,
	presentation_id: str,
	slide_id: int,
	shape_ids: list[int],
) -> dict:
	"""Remove one or many shapes by shape IDs from a specific slide."""
	if presentation_id not in presentations:
		return {"error": "Presentation ID not found."}

	presentation_file = presentations[presentation_id]
	pptx_object = presentation_file.get_pptx_object()

	slide_to_update = get_slide_by_id(pptx_object, slide_id)
	if slide_to_update is None:
		return {"error": "Slide ID not found."}

	if not isinstance(shape_ids, list) or len(shape_ids) == 0:
		return {"error": "shape_ids must be a non-empty list."}

	shape_ids_int: list[int] = []
	invalid_shape_ids: list = []
	for raw_shape_id in shape_ids:
		try:
			shape_ids_int.append(int(raw_shape_id))
		except (TypeError, ValueError):
			invalid_shape_ids.append(raw_shape_id)

	if invalid_shape_ids:
		return {
			"error": "All shape_ids must be numeric values.",
			"invalid_shape_ids": invalid_shape_ids,
		}

	requested_shape_ids = list(dict.fromkeys(shape_ids_int))
	shape_map = {shape.shape_id: shape for shape in slide_to_update.shapes}

	removed_shapes: list[dict] = []
	not_found_shape_ids: list[int] = []

	for shape_id_int in requested_shape_ids:
		shape_to_remove = shape_map.get(shape_id_int)
		if shape_to_remove is None:
			not_found_shape_ids.append(shape_id_int)
			continue

		shape_element = shape_to_remove.element
		shape_parent = shape_element.getparent()
		if shape_parent is None:
			not_found_shape_ids.append(shape_id_int)
			continue

		removed_shapes.append(
			{
				"shape_id": shape_to_remove.shape_id,
				"name": shape_to_remove.name,
				"shape_type": str(shape_to_remove.shape_type),
				"shape_type_value": int(shape_to_remove.shape_type),
			}
		)
		shape_parent.remove(shape_element)

	if not removed_shapes:
		return {
			"error": "No matching shapes were removed.",
			"presentation_id": presentation_id,
			"slide_id": slide_id,
			"requested_shape_ids": requested_shape_ids,
			"not_found_shape_ids": not_found_shape_ids,
		}

	return {
		"message": "Shapes removed from slide successfully",
		"presentation_id": presentation_id,
		"slide_id": slide_id,
		"requested_shape_ids": requested_shape_ids,
		"removed_shapes": removed_shapes,
		"not_found_shape_ids": not_found_shape_ids,
	}


def extract_text_from_shape(shape) -> str | None:
	"""Safely extract visible text from a shape."""
	if hasattr(shape, "text") and shape.text:
		return shape.text
	if hasattr(shape, "text_frame") and shape.text_frame is not None:
		return shape.text_frame.text
	return None


def serialize_placeholder(placeholder) -> dict:
	"""Convert a placeholder into a serializable dictionary."""
	return {
		"placeholder_shape_id": placeholder.shape_id,
		"placeholder_type": str(placeholder.placeholder_format.type),
		"placeholder_type_value": int(placeholder.placeholder_format.type),
		"placeholder_text": extract_text_from_shape(placeholder),
	}


def serialize_shape(shape) -> dict:
	"""Convert a shape into a serializable dictionary."""
	return {
		"shape_id": shape.shape_id,
		"name": shape.name,
		"shape_type": str(shape.shape_type),
		"shape_type_value": int(shape.shape_type),
		"left": int(shape.left),
		"top": int(shape.top),
		"width": int(shape.width),
		"height": int(shape.height),
		"has_text": hasattr(shape, "text_frame") and shape.text_frame is not None,
		"shape_text": extract_text_from_shape(shape),
	}


def serialize_slide(slide, slide_index: int) -> dict:
	"""Convert a slide into a serializable dictionary."""
	notes_text = None
	if slide.has_notes_slide and slide.notes_slide is not None:
		notes_text = slide.notes_slide.notes_text_frame.text

	placeholders = [serialize_placeholder(p) for p in slide.placeholders]
	shapes = [serialize_shape(s) for s in slide.shapes]

	return {
		"slide_id": slide.slide_id,
		"slide_index": slide_index,
		"name": slide.name,
		"slide_layout": slide.slide_layout.name,
		"has_notes": slide.has_notes_slide,
		"slide_notes": notes_text,
		"placeholders_count": len(placeholders),
		"shapes_count": len(shapes),
		"slide_placeholders": placeholders,
		"slide_shapes": shapes,
	}


def _to_points(emu_value: int) -> float:
	"""Convert EMU value to points."""
	return max(0.0, float(emu_value) / EMU_PER_POINT)


def _font_size_pt_from_oxml_rpr(rpr) -> Optional[float]:
	"""Read OOXML font size from an rPr-like node (sz is in 1/100 pt)."""
	if rpr is None:
		return None

	sz_attr = rpr.get("sz")
	if sz_attr is None:
		return None

	try:
		return float(int(sz_attr)) / 100.0
	except (TypeError, ValueError):
		return None


def _font_size_pt_from_paragraph_xml(paragraph) -> Optional[float]:
	"""Read font size directly from paragraph OOXML when available."""
	p_element = getattr(paragraph, "_p", None)
	if p_element is None:
		return None

	# 1) End paragraph run properties can carry inherited effective size.
	end_para_rpr = p_element.find(qn("a:endParaRPr"))
	font_size_pt = _font_size_pt_from_oxml_rpr(end_para_rpr)
	if font_size_pt is not None:
		return font_size_pt

	# 2) Paragraph properties may define default run properties.
	p_pr = p_element.find(qn("a:pPr"))
	if p_pr is not None:
		def_rpr = p_pr.find(qn("a:defRPr"))
		font_size_pt = _font_size_pt_from_oxml_rpr(def_rpr)
		if font_size_pt is not None:
			return font_size_pt

	return None


def _font_size_pt_from_list_style(paragraph, shape) -> Optional[float]:
	"""Resolve paragraph-level default font size from text body list style."""
	if shape is None or not hasattr(shape, "text_frame") or shape.text_frame is None:
		return None

	tx_body = getattr(shape.text_frame, "_txBody", None)
	if tx_body is None:
		return None

	lst_style = tx_body.find(qn("a:lstStyle"))
	if lst_style is None:
		return None

	# Paragraph level is 0-based; OOXML list style tags are lvl1..lvl9.
	level = max(0, min(8, int(getattr(paragraph, "level", 0))))
	level_tag = qn(f"a:lvl{level + 1}pPr")
	level_p_pr = lst_style.find(level_tag)
	if level_p_pr is None:
		return None

	def_rpr = level_p_pr.find(qn("a:defRPr"))
	return _font_size_pt_from_oxml_rpr(def_rpr)


def _font_size_pt_from_layout_placeholder(paragraph, shape) -> Optional[float]:
	"""Resolve font size from the matching layout placeholder list style."""
	if shape is None:
		return None

	try:
		placeholder_idx = int(shape.placeholder_format.idx)
	except Exception:
		return None

	slide_layout = getattr(getattr(shape, "part", None), "slide_layout", None)
	if slide_layout is None:
		return None

	for layout_placeholder in slide_layout.placeholders:
		try:
			if int(layout_placeholder.placeholder_format.idx) != placeholder_idx:
				continue
		except Exception:
			continue

		font_size_pt = _font_size_pt_from_paragraph_xml(layout_placeholder.text_frame.paragraphs[0])
		if font_size_pt is not None:
			return font_size_pt

		font_size_pt = _font_size_pt_from_list_style(paragraph, layout_placeholder)
		if font_size_pt is not None:
			return font_size_pt

		break

	return None


def _get_paragraph_font_size_pt(paragraph, shape=None) -> float:
	"""Resolve paragraph font size in points with inherited-style fallback."""
	if paragraph.font is not None and paragraph.font.size is not None:
		return float(paragraph.font.size.pt)

	for run in paragraph.runs:
		if run.font is not None and run.font.size is not None:
			return float(run.font.size.pt)

	font_size_pt = _font_size_pt_from_paragraph_xml(paragraph)
	if font_size_pt is not None:
		return font_size_pt

	font_size_pt = _font_size_pt_from_list_style(paragraph, shape)
	if font_size_pt is not None:
		return font_size_pt

	font_size_pt = _font_size_pt_from_layout_placeholder(paragraph, shape)
	if font_size_pt is not None:
		return font_size_pt

	return DEFAULT_FONT_SIZE_PT


def analyze_text_overflow_in_shape(shape) -> dict:
	"""Estimate whether text likely overflows the shape's text area."""
	if not hasattr(shape, "text_frame") or shape.text_frame is None:
		return {
			"overflow_detected": False,
			"reason": "shape_has_no_text_frame",
		}

	text_frame = shape.text_frame
	available_width_emu = int(shape.width) - int(text_frame.margin_left) - int(text_frame.margin_right)
	available_height_emu = int(shape.height) - int(text_frame.margin_top) - int(text_frame.margin_bottom)

	if available_width_emu <= 0 or available_height_emu <= 0:
		return {
			"overflow_detected": True,
			"reason": "invalid_text_area",
			"available_width_pt": _to_points(available_width_emu),
			"available_height_pt": _to_points(available_height_emu),
		}

	available_width_pt = _to_points(available_width_emu)
	available_height_pt = _to_points(available_height_emu)
	word_wrap = True if text_frame.word_wrap is None else bool(text_frame.word_wrap)

	total_required_height_pt = 0.0
	longest_estimated_line_pt = 0.0
	horizontal_overflow_detected = False

	for paragraph in text_frame.paragraphs:
		paragraph_text = paragraph.text or ""
		font_size_pt = _get_paragraph_font_size_pt(paragraph, shape)
		line_height_pt = font_size_pt * LINE_HEIGHT_FACTOR
		chars_per_line = max(1, int(available_width_pt / (font_size_pt * AVG_CHAR_WIDTH_FACTOR)))

		raw_lines = paragraph_text.splitlines() or [""]
		for raw_line in raw_lines:
			line_char_count = len(raw_line)
			estimated_line_width_pt = line_char_count * font_size_pt * AVG_CHAR_WIDTH_FACTOR
			longest_estimated_line_pt = max(longest_estimated_line_pt, estimated_line_width_pt)

			if word_wrap:
				wrapped_line_count = max(1, math.ceil(line_char_count / chars_per_line))
				total_required_height_pt += wrapped_line_count * line_height_pt
			else:
				total_required_height_pt += line_height_pt
				if estimated_line_width_pt > available_width_pt:
					horizontal_overflow_detected = True

	vertical_overflow_detected = total_required_height_pt > available_height_pt
	overflow_detected = horizontal_overflow_detected or vertical_overflow_detected

	return {
		"overflow_detected": overflow_detected,
		"horizontal_overflow_detected": horizontal_overflow_detected,
		"vertical_overflow_detected": vertical_overflow_detected,
		"word_wrap": word_wrap,
		"available_width_pt": round(available_width_pt, 2),
		"available_height_pt": round(available_height_pt, 2),
		"required_height_pt_estimate": round(total_required_height_pt, 2),
		"longest_line_width_pt_estimate": round(longest_estimated_line_pt, 2),
	}
