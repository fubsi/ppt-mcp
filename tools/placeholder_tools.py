from mcp.types import ToolAnnotations
from pptx.enum.shapes import PP_PLACEHOLDER

from utils.helper_methods import (
    analyze_text_overflow_in_shape,
    get_placeholder_by_shape_id,
    get_slide_by_id,
    resolve_picture_source,
)

def register_placeholder_tools(app, presentations):
    """Register placeholder tools for PowerPoint file management."""

    @app.tool(
        annotations=ToolAnnotations(
            title="Insert picture into placeholder",
            readOnlyHint=False,
        ),
        description="Inserts an image into a PICTURE placeholder on a specific slide.\
            The presentation_id, slide_id, placeholder_shape_id, and image_path must be provided for this operation.\
            image_path can be a local file path or an http/https URL.\
            Use get_slides to find slide_id and placeholder_shape_id values.\
            Placeholder-first policy: this is the preferred way to add content before manual shape tools."
    )
    def insert_picture_into_placeholder(
        presentation_id: str,
        slide_id: int,
        placeholder_shape_id: int,
        image_path: str,
    ) -> dict:
        """Insert an image from a local path or URL into a picture placeholder."""
        if presentation_id not in presentations:
            return {"error": "Presentation ID not found."}

        image_source, image_source_label, source_error = resolve_picture_source(image_path)
        if source_error is not None:
            return source_error

        presentation_file = presentations[presentation_id]
        pptx_object = presentation_file.get_pptx_object()

        slide_to_update = get_slide_by_id(pptx_object, slide_id)

        if slide_to_update is None:
            return {"error": "Slide ID not found."}

        target_placeholder = get_placeholder_by_shape_id(slide_to_update, placeholder_shape_id)

        if target_placeholder is None:
            return {"error": "Placeholder shape ID not found on this slide."}

        placeholder_type = target_placeholder.placeholder_format.type
        if int(placeholder_type) != int(PP_PLACEHOLDER.PICTURE):
            return {
                "error": "Target placeholder is not a PICTURE placeholder.",
                "placeholder_type": str(placeholder_type),
                "placeholder_type_value": int(placeholder_type),
            }

        inserted_picture = target_placeholder.insert_picture(image_source)

        return {
            "message": "Picture inserted into placeholder successfully",
            "presentation_id": presentation_id,
            "slide_id": slide_id,
            "placeholder_shape_id": placeholder_shape_id,
            "image_path": image_source_label,
            "inserted_picture": {
                "shape_id": inserted_picture.shape_id,
                "name": inserted_picture.name,
                "shape_type": str(inserted_picture.shape_type),
                "left": int(inserted_picture.left),
                "top": int(inserted_picture.top),
                "width": int(inserted_picture.width),
                "height": int(inserted_picture.height),
            },
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Insert text into placeholder",
            readOnlyHint=False,
        ),
        description="Inserts text into a text-capable placeholder on a specific slide.\
            The presentation_id, slide_id, placeholder_shape_id, and text must be provided for this operation.\
            Use get_slides to find slide_id and placeholder_shape_id values.\
            Placeholder-first policy: this is the preferred way to add content before manual shape tools."
    )
    def insert_text_into_placeholder(
        presentation_id: str,
        slide_id: int,
        placeholder_shape_id: int,
        text: str,
    ) -> dict:
        """Insert text into an existing text-capable placeholder."""
        if presentation_id not in presentations:
            return {"error": "Presentation ID not found."}

        presentation_file = presentations[presentation_id]
        pptx_object = presentation_file.get_pptx_object()

        slide_to_update = get_slide_by_id(pptx_object, slide_id)

        if slide_to_update is None:
            return {"error": "Slide ID not found."}

        target_placeholder = get_placeholder_by_shape_id(slide_to_update, placeholder_shape_id)

        if target_placeholder is None:
            return {"error": "Placeholder shape ID not found on this slide."}

        placeholder_type = target_placeholder.placeholder_format.type
        if not hasattr(target_placeholder, "text_frame") or target_placeholder.text_frame is None:
            return {
                "error": "Target placeholder does not support text insertion.",
                "placeholder_type": str(placeholder_type),
                "placeholder_type_value": int(placeholder_type),
            }

        target_placeholder.text_frame.clear()
        target_placeholder.text_frame.text = text
        overflow_analysis = analyze_text_overflow_in_shape(target_placeholder)

        return {
            "message": "Text inserted into placeholder successfully",
            "presentation_id": presentation_id,
            "slide_id": slide_id,
            "placeholder_shape_id": placeholder_shape_id,
            "text": target_placeholder.text_frame.text,
            "text_overflow_detected": overflow_analysis["overflow_detected"],
            "text_overflow_analysis": overflow_analysis,
            "placeholder_type": str(placeholder_type),
            "placeholder_type_value": int(placeholder_type),
        }